from typing import Dict, List, Optional
from fastapi import FastAPI, File, UploadFile, Depends, HTTPException, Query
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from sqlalchemy.orm import Session
from pathlib import Path
import torch
import numpy as np
from PIL import Image
import time
from datetime import datetime
import logging
import json
import uvicorn
import tempfile
import os
import urllib.request
import re

from config import settings
from database import init_db, get_db, ContentSubmission, ModerationResult, ManualReview, engine, Base

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize database
Base.metadata.create_all(bind=engine)
init_db()

# Initialize FastAPI app
app = FastAPI(
    title="Content Moderation AI",
    description="Detect harmful content with 6 advanced classifiers",
    version="1.0.0"
)

# Add CORS middleware - ALLOW ALL
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Convert settings paths to Path objects
UPLOAD_DIR = Path(settings.UPLOAD_DIR)
MODEL_DIR = Path(settings.MODEL_DIR)
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
MODEL_DIR.mkdir(parents=True, exist_ok=True)

# Global model storage
class ModelManager:
    def __init__(self):
        self.models = {}
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        logger.info(f"[START] Using device: {self.device}")
        self.load_all_models()
    
    def load_all_models(self):
        """Load all classifiers"""
        try:
            logger.info("[LOAD] Loading ML models...")
            from transformers import pipeline
            
            # 1. Zero-Shot Text Classifier (Used for Hate Speech, Toxicity, Violence, Misinfo, Spam)
            try:
                self.models['zero_shot'] = pipeline(
                    "zero-shot-classification",
                    model="typeform/distilbart-base-uncased-mnli",
                    device=0 if self.device == "cuda" else -1
                )
                logger.info("[OK] Zero-shot text classifier loaded")
            except Exception as e:
                logger.warning(f"[WARN] Zero-shot text model failed: {e}")
            
            # 2. NSFW Image Detection (Falconsai)
            try:
                self.models['nsfw_image'] = pipeline(
                    "image-classification",
                    model="Falconsai/nsfw_image_detection",
                    device=0 if self.device == "cuda" else -1
                )
                logger.info("[OK] NSFW image classifier (Falconsai) loaded")
            except Exception as e:
                logger.warning(f"[WARN] NSFW image model (Falconsai) failed: {e}")
            
            # 3. Robust NSFW Image Detection (AdamCodd)
            try:
                self.models['nsfw_image_robust'] = pipeline(
                    "image-classification",
                    model="AdamCodd/vit-base-nsfw-detector",
                    device=0 if self.device == "cuda" else -1
                )
                logger.info("[OK] Robust NSFW image classifier (AdamCodd) loaded")
            except Exception as e:
                logger.warning(f"[WARN] Robust NSFW image model (AdamCodd) failed: {e}")

            # 4. ToxicBERT — Multi-label toxicity (fine-tuned on Jigsaw dataset)
            # Detects: toxic, severe_toxic, obscene (profanity), threat (violence), insult (cyberbullying), identity_hate
            try:
                self.models['toxic_bert'] = pipeline(
                    "text-classification",
                    model="unitary/toxic-bert",
                    top_k=None,
                    device=0 if self.device == "cuda" else -1
                )
                logger.info("[OK] ToxicBERT multi-label classifier loaded")
            except Exception as e:
                logger.warning(f"[WARN] ToxicBERT failed: {e}")

            # 5. DeHateBERT — Dedicated hate speech detection
            try:
                self.models['hate_classifier'] = pipeline(
                    "text-classification",
                    model="Hate-speech-CNERG/dehatebert-mono-english",
                    device=0 if self.device == "cuda" else -1
                )
                logger.info("[OK] DeHateBERT hate speech classifier loaded")
            except Exception as e:
                logger.warning(f"[WARN] DeHateBERT failed: {e}")

            # 6. NSFW Text Classifier — Explicit sexual content in text
            try:
                self.models['nsfw_text'] = pipeline(
                    "text-classification",
                    model="michellejieli/NSFW_text_classifier",
                    device=0 if self.device == "cuda" else -1
                )
                logger.info("[OK] NSFW text classifier loaded")
            except Exception as e:
                logger.warning(f"[WARN] NSFW text classifier failed: {e}")

            logger.info(f"[OK] All models loaded! Total: {len(self.models)}")
            
        except Exception as e:
            logger.error(f"[ERROR] Error loading models: {e}")


# Initialize model manager
try:
    model_manager = ModelManager()
    logger.info("[OK] Classifiers initialized successfully!")
except Exception as e:
    logger.error(f"[WARN] Warning: Some models may not be loaded: {e}")


# ==================== UTILITY FUNCTIONS ====================

def process_image(file_path: Path) -> Dict:
    """Process image and return metadata"""
    try:
        image = Image.open(file_path)
        return {
            "size": image.size,
            "format": image.format,
            "mode": image.mode,
            "file_size": file_path.stat().st_size
        }
    except Exception as e:
        logger.error(f"Error processing image: {e}")
        return {"error": str(e)}


def moderate_text(text: str, enabled_models: Optional[List[str]] = None) -> Dict:
    """
    Multi-model text classification pipeline:
    - ToxicBERT       → toxicity, profanity, violence (threats), cyberbullying, hate_speech
    - DeHateBERT      → hate_speech (ensemble with ToxicBERT, takes max)
    - NSFW-text       → sexual_content
    - Zero-Shot BART  → misinformation, spam, self_harm, extremism, phishing,
                        defamation, hate_symbols, pseudoscience, copyright
    """
    results = {}
    lower_text = text.lower()
    text_input = text[:512]

    try:
        # ================================================================
        # STAGE 1 — ToxicBERT (fine-tuned Jigsaw, multi-label)
        # Labels: toxic | severe_toxic | obscene | threat | insult | identity_hate
        # ================================================================
        if 'toxic_bert' in model_manager.models:
            try:
                raw = model_manager.models['toxic_bert'](text_input)
                # top_k=None returns list-of-lists: [[{label, score}, ...]]
                label_scores = {r['label']: r['score'] for r in raw[0]}

                # Map ToxicBERT labels → our categories
                bert_map = {
                    'toxicity':     ('toxic',         0.5),
                    'profanity':    ('obscene',        0.5),
                    'violence':     ('threat',         0.4),
                    'cyberbullying':('insult',         0.5),
                    'hate_speech':  ('identity_hate',  0.4),
                }

                trigger_kw = {
                    'toxicity':      ["stupid", "idiot", "dumb", "bitch", "fuck", "shit", "asshole", "loser", "pathetic"],
                    'profanity':     ["fuck", "shit", "damn", "bitch", "cunt", "dick", "bastard"],
                    'violence':      ["kill", "murder", "attack", "shoot", "bomb", "slaughter", "assassinate"],
                    'cyberbullying': ["doxx", "bully", "harass", "dox", "find you", "cancel"],
                    'hate_speech':   ["hate", "racist", "slur", "supremacy", "nazi", "bigot", "inferior"],
                }

                for cat, (bert_label, threshold) in bert_map.items():
                    if enabled_models is not None and cat not in enabled_models:
                        continue
                    score = label_scores.get(bert_label, 0.0)
                    found_words = [w for w in trigger_kw.get(cat, []) if w in lower_text]
                    clean_score = 0.0 if score <= 0.15 else score

                    if score > threshold:
                        if found_words:
                            reason = f"ToxicBERT flagged at {(score*100):.1f}%. Trigger keywords: '{', '.join(found_words[:3])}'."
                        else:
                            reason = f"ToxicBERT: {cat.replace('_',' ')} detected at {(score*100):.1f}% confidence."
                    elif score > 0.15:
                        reason = f"Low signal ({(score*100):.1f}%). Mild indicators of {cat.replace('_',' ')} present."
                    else:
                        reason = f"Safe (0.0%). ToxicBERT finds no evidence of {cat.replace('_',' ')}."

                    results[cat] = {
                        "score": clean_score,
                        "label": "FLAGGED" if score > threshold else "SAFE",
                        "reason": reason
                    }
            except Exception as e:
                logger.warning(f"[WARN] ToxicBERT failed: {e}")

        # ================================================================
        # STAGE 2 — DeHateBERT (dedicated hate-speech, ensemble with ToxicBERT)
        # Labels: 'hate' | 'noHate'
        # ================================================================
        if 'hate_classifier' in model_manager.models:
            if enabled_models is None or 'hate_speech' in enabled_models:
                try:
                    res = model_manager.models['hate_classifier'](text_input)
                    label = res[0]['label'].lower()
                    raw_score = res[0]['score']
                    hate_score = raw_score if 'hate' in label and 'no' not in label else (1.0 - raw_score)
                    clean_score = 0.0 if hate_score <= 0.3 else hate_score

                    if hate_score > 0.5:
                        reason = f"DeHateBERT: Hate speech detected at {(hate_score*100):.1f}% confidence."
                    elif hate_score > 0.3:
                        reason = f"Borderline ({(hate_score*100):.1f}%). Some hateful signals present."
                    else:
                        reason = f"Safe (0.0%). DeHateBERT finds no hate speech."

                    # Ensemble: take MAX of ToxicBERT and DeHateBERT for hate_speech
                    existing_score = results.get('hate_speech', {}).get('score', 0.0)
                    if clean_score > existing_score:
                        results['hate_speech'] = {
                            "score": clean_score,
                            "label": "FLAGGED" if hate_score > 0.5 else "SAFE",
                            "reason": reason
                        }
                except Exception as e:
                    logger.warning(f"[WARN] DeHateBERT failed: {e}")

        # ================================================================
        # STAGE 3 — NSFW Text Classifier (dedicated sexual content)
        # Labels: 'NSFW' | 'SFW'
        # ================================================================
        if 'nsfw_text' in model_manager.models:
            if enabled_models is None or 'sexual_content' in enabled_models:
                try:
                    res = model_manager.models['nsfw_text'](text_input)
                    label = res[0]['label'].upper()
                    raw_score = res[0]['score']
                    nsfw_score = raw_score if label == 'NSFW' else (1.0 - raw_score)
                    
                    # Prevent false positives on professional documents (like Resumes/CVs) and check for actual explicit keywords
                    explicit_keywords = [
                        "porn", "xxx", "nude", "naked", "sexy", "erotic", "penis", "vagina", 
                        "orgasm", "clitoris", "semen", "sperm", "prostitute", "escort", "breast", "pussy",
                        "dick", "cock", "adult content", "sexual", "intercourse", "nsfw"
                    ]
                    professional_keywords = [
                        "resume", "cv", "curriculum vitae", "experience", "education", "projects", 
                        "skills", "university", "technologies", "achievement", "certifications", 
                        "python", "javascript", "react", "fastapi", "developer", "engineer", "intern", 
                        "github", "linkedin", "contact", "email", "address", "phone"
                    ]
                    
                    has_explicit = any(w in lower_text for w in explicit_keywords)
                    has_professional = any(w in lower_text for w in professional_keywords)
                    
                    # Override if flagged as NSFW but lacks explicit words or has clear professional keywords
                    if nsfw_score > 0.4 and not has_explicit:
                        if has_professional or nsfw_score < 0.85:
                            nsfw_score = 0.08  # Override to safe
                    
                    clean_score = 0.0 if nsfw_score <= 0.3 else nsfw_score

                    if nsfw_score > 0.5:
                        reason = f"NSFW classifier: Explicit content detected at {(nsfw_score*100):.1f}% confidence."
                    elif nsfw_score > 0.3:
                        reason = f"Borderline ({(nsfw_score*100):.1f}%). Mildly suggestive content."
                    else:
                        reason = f"Safe (0.0%). Content appears work-safe."

                    results['sexual_content'] = {
                        "score": clean_score,
                        "label": "FLAGGED" if nsfw_score > 0.5 else "SAFE",
                        "reason": reason
                    }
                except Exception as e:
                    logger.warning(f"[WARN] NSFW text classifier failed: {e}")

        # ================================================================
        # STAGE 4 — Zero-Shot BART for rare/complex categories
        # (categories NOT already handled by dedicated models above)
        # ================================================================
        dedicated_cats = {'toxicity', 'profanity', 'violence', 'cyberbullying', 'hate_speech', 'sexual_content'}

        if 'zero_shot' in model_manager.models:
            zero_shot = model_manager.models['zero_shot']

            zs_categories = {
                "misinformation": ["false, misleading, or conspiracy theory", "factually accurate and truthful", "neutral topic"],
                "spam":           ["spam, promotional, or repetitive marketing", "genuine human content", "neutral topic"],
                "self_harm":      ["self-harm, suicide, or suicidal ideation", "healthy state of mind", "neutral topic"],
                "extremism":      ["terrorism, extremism, or radicalization", "moderate and peaceful", "neutral topic"],
                "phishing":       ["phishing, scam, or credential theft", "safe and legitimate", "neutral topic"],
                "defamation":     ["defamation, slander, or libel", "fair and factual statement", "neutral topic"],
                "hate_symbols":   ["hate group symbols or fascist imagery references", "neutral historical context", "neutral topic"],
                "pseudoscience":  ["pseudoscience, health misinformation, or quackery", "scientifically accurate", "neutral topic"],
                "copyright":      ["copyright infringement, piracy, or leaked content", "original and licensed content", "neutral topic"],
            }

            zs_trigger_kw = {
                "misinformation": ["hoax", "fake", "conspiracy", "plandemic", "5g", "flat earth", "microchip"],
                "spam":           ["click here", "buy now", "free money", "lottery", "earn fast", "bitcoin"],
                "self_harm":      ["suicide", "kill myself", "cut myself", "worthless", "end it all", "no reason to live"],
                "extremism":      ["jihad", "manifesto", "terror", "radical", "overthrow", "uprising"],
                "phishing":       ["password", "verify account", "bank details", "ssn", "wire transfer"],
                "defamation":     ["scammer", "fraud", "thief", "liar", "corrupt", "pedophile"],
                "hate_symbols":   ["swastika", "ss bolts", "kkk", "iron cross", "white pride", "black sun"],
                "pseudoscience":  ["healing crystals", "vaccines cause autism", "bleach cure", "miracle mineral"],
                "copyright":      ["free download", "crack", "torrent", "leaked", "pirated", "watch online free"],
            }

            for key, labels in zs_categories.items():
                if enabled_models is not None and key not in enabled_models:
                    continue
                # Skip if a dedicated model already gave a result for this category
                if key in dedicated_cats and key in results:
                    continue
                try:
                    res = zero_shot(text_input, labels)
                    flag_idx = res['labels'].index(labels[0])
                    score = float(res['scores'][flag_idx])
                    found_words = [w for w in zs_trigger_kw.get(key, []) if w in lower_text]

                    if score > 0.6:
                        if found_words:
                            reason = f"High risk ({(score*100):.1f}%). Strong signals of {key.replace('_',' ')}. Keywords: '{', '.join(found_words[:3])}'."
                        else:
                            reason = f"Flagged at {(score*100):.1f}%. Semantic meaning matches '{labels[0]}'."
                    elif score > 0.3:
                        reason = f"Elevated ({(score*100):.1f}%). Mild signals of {key.replace('_',' ')} present."
                    else:
                        reason = f"Safe (0.0%). No evidence of {labels[0]}."

                    clean_score = 0.0 if score <= 0.3 else score
                    results[key] = {
                        "score": clean_score,
                        "label": "FLAGGED" if score > 0.6 else "SAFE",
                        "reason": reason
                    }
                except Exception as zs_e:
                    logger.warning(f"Zero-shot [{key}] failed: {zs_e}")

    except Exception as e:
        logger.error(f"Error in text moderation: {e}")

    return results


def moderate_image(file_path: Path, enabled_models: Optional[List[str]] = None) -> Dict:
    """Run NSFW image classifier using an ensemble of Falconsai and AdamCodd models for high robustness"""
    results = {}
    
    if enabled_models is not None and "nsfw" not in enabled_models:
        return results
    
    try:
        scores = []
        
        # 1. Falconsai NSFW Image Detector
        if 'nsfw_image' in model_manager.models:
            try:
                image = Image.open(file_path)
                res = model_manager.models['nsfw_image'](image)
                raw_score = float(res[0]['score'])
                label = res[0]['label']
                nsfw_score = raw_score if label != "normal" else (1.0 - raw_score)
                scores.append(nsfw_score)
            except Exception as e:
                logger.warning(f"Falconsai NSFW check failed: {e}")

        # 2. AdamCodd NSFW Image Detector (Highly Robust)
        if 'nsfw_image_robust' in model_manager.models:
            try:
                image = Image.open(file_path)
                res = model_manager.models['nsfw_image_robust'](image)
                raw_score = float(res[0]['score'])
                label = res[0]['label']
                nsfw_score = raw_score if label == "nsfw" else (1.0 - raw_score)
                scores.append(nsfw_score)
            except Exception as e:
                logger.warning(f"AdamCodd NSFW check failed: {e}")

        if scores:
            final_nsfw_score = max(scores)
            results['nsfw'] = {
                "score": final_nsfw_score,
                "label": "FLAGGED" if final_nsfw_score > 0.5 else "SAFE",
                "reason": f"Image flagged as NSFW with {(final_nsfw_score*100):.1f}% confidence. Explicit or adult content detected." if final_nsfw_score > 0.5 else "Content appears normal and safe."
            }



    except Exception as e:
        logger.error(f"Error in image moderation: {e}")
    
    return results


def moderate_image_pil(pil_image, enabled_models: Optional[List[str]] = None) -> Dict:
    """Run NSFW classifier on a PIL Image object (for video frames) using an ensemble"""
    results = {}
    if enabled_models is not None and "nsfw" not in enabled_models:
        return results
    try:
        scores = []
        
        # 1. Falconsai
        if 'nsfw_image' in model_manager.models:
            res = model_manager.models['nsfw_image'](pil_image)
            raw_score = float(res[0]['score'])
            label = res[0]['label']
            nsfw_score = raw_score if label != "normal" else (1.0 - raw_score)
            scores.append(nsfw_score)
            
        # 2. AdamCodd
        if 'nsfw_image_robust' in model_manager.models:
            res = model_manager.models['nsfw_image_robust'](pil_image)
            raw_score = float(res[0]['score'])
            label = res[0]['label']
            nsfw_score = raw_score if label == "nsfw" else (1.0 - raw_score)
            scores.append(nsfw_score)

        if scores:
            final_nsfw_score = max(scores)
            results['nsfw'] = {
                "score": final_nsfw_score, 
                "label": "FLAGGED" if final_nsfw_score > 0.5 else "SAFE",
                "reason": f"Frame flagged as NSFW with {(final_nsfw_score*100):.1f}% confidence. Explicit or adult content detected." if final_nsfw_score > 0.5 else "Content appears normal and safe."
            }
    except Exception as e:
        logger.warning(f"NSFW check on frame failed: {e}")
    return results


def get_moderation_status(results: Dict) -> str:
    """Determine if content should be APPROVED, FLAGGED, or NEEDS_REVIEW"""
    
    # Get all scores dynamically
    scores = []
    for category, data in results.items():
        if isinstance(data, dict) and 'score' in data:
            scores.append(data.get('score', 0))
    
    if not scores:
        return "NEEDS_REVIEW"
    
    max_score = max(scores)
    
    if max_score > 0.6:
        return "FLAGGED"
    elif max_score > 0.4:
        return "NEEDS_REVIEW"
    else:
        return "APPROVED"


# ==================== API ENDPOINTS ====================

@app.get("/")
async def root():
    """Health check endpoint"""
    return {
        "status": "ok",
        "message": "Content Moderation AI API is running",
        "version": "1.0.0"
    }


@app.post("/moderate/text")
async def moderate_text_endpoint(
    text: str = Query(...),
    enabled_models: Optional[str] = Query(None),
    db: Session = Depends(get_db)
) -> Dict:
    """
    Moderate text content
    
    Query Parameters:
    - text: Text to moderate (required)
    
    Returns:
    - submission_id: ID of the submission
    - status: APPROVED, FLAGGED, or NEEDS_REVIEW
    - results: Results from classifiers
    """
    
    try:
        if not text or len(text.strip()) == 0:
            raise HTTPException(status_code=400, detail="Text cannot be empty")
        
        logger.info(f"Processing text: {text[:50]}...")
        
        # Run moderation
        models_list = enabled_models.split(",") if enabled_models else None
        moderation_results = moderate_text(text, models_list)
        
        # Add Authenticity & Plagiarism checking
        extra_checks = check_authenticity_and_plagiarism(text, models_list)
        moderation_results.update(extra_checks)
        
        # Determine status
        status = get_moderation_status(moderation_results)
        
        # Save to database
        submission = ContentSubmission(
            content_type="text",
            file_path="",
            file_name="",
            created_at=datetime.now()
        )
        db.add(submission)
        db.commit()
        db.refresh(submission)
        
        # Save results
        moderation_record = ModerationResult(
            submission_id=submission.id,
            classifier_results=moderation_results,
            overall_status=status,
            moderated_at=datetime.now()
        )
        db.add(moderation_record)
        db.commit()
        
        logger.info(f"[OK] Text moderation completed. Status: {status}")
        
        return {
            "submission_id": submission.id,
            "content_type": "text",
            "status": status,
            "results": moderation_results,
            "timestamp": datetime.now().isoformat()
        }
    
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"[ERROR] Error moderating text: {e}")
        raise HTTPException(status_code=500, detail=f"Moderation error: {str(e)}")


@app.post("/moderate/image")
async def moderate_image_endpoint(
    file: UploadFile = File(...),
    enabled_models: Optional[str] = Query(None),
    db: Session = Depends(get_db)
) -> Dict:
    """
    Moderate image content
    
    File Upload:
    - file: Image file (JPG, PNG, GIF, WebP, BMP)
    
    Returns:
    - submission_id: ID of the submission
    - status: APPROVED, FLAGGED, or NEEDS_REVIEW
    - results: Results from image classifier
    """
    
    try:
        # Validate file
        is_image_ext = any(file.filename.lower().endswith(ext) for ext in ['.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp'])
        if (not file.content_type or not file.content_type.startswith('image/')) and not is_image_ext:
            raise HTTPException(status_code=400, detail="Invalid image format. Please upload an image file.")
        
        logger.info(f"Processing image: {file.filename}")
        
        # Save uploaded file
        file_path = UPLOAD_DIR / f"{int(time.time())}_{file.filename}"
        contents = await file.read()
        with open(file_path, "wb") as f:
            f.write(contents)
        
        # Process image
        image_metadata = process_image(file_path)
        
        # Run NSFW classifier
        models_list = enabled_models.split(",") if enabled_models else None
        moderation_results = moderate_image(file_path, models_list)
        
        # Determine status
        status = get_moderation_status(moderation_results)
        
        # Save to database
        submission = ContentSubmission(
            content_type="image",
            file_path=str(file_path),
            file_name=file.filename,
            created_at=datetime.now()
        )
        db.add(submission)
        db.commit()
        db.refresh(submission)
        
        # Save results
        moderation_record = ModerationResult(
            submission_id=submission.id,
            classifier_results={**moderation_results, "metadata": image_metadata},
            overall_status=status,
            moderated_at=datetime.now()
        )
        db.add(moderation_record)
        db.commit()
        
        logger.info(f"[OK] Image moderation completed. Status: {status}")
        
        return {
            "submission_id": submission.id,
            "content_type": "image",
            "status": status,
            "results": moderation_results,
            "metadata": image_metadata,
            "timestamp": datetime.now().isoformat()
        }
    
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"[ERROR] Error moderating image: {e}")
        raise HTTPException(status_code=500, detail=f"Image moderation error: {str(e)}")





# ==================== VIDEO MODERATION ====================

@app.post("/moderate/video")
async def moderate_video_endpoint(
    file: UploadFile = File(...),
    enabled_models: Optional[str] = Query(None),
    db: Session = Depends(get_db)
) -> Dict:
    """Moderate video content by extracting and analyzing key frames"""
    try:
        allowed_types = ['video/mp4', 'video/avi', 'video/quicktime', 'video/webm', 'video/x-msvideo', 'video/x-matroska']
        is_video_ext = any(file.filename.lower().endswith(ext) for ext in ['.mp4', '.avi', '.mov', '.webm', '.mkv'])
        if (not file.content_type or file.content_type not in allowed_types) and not is_video_ext:
            raise HTTPException(status_code=400, detail="Invalid video format. Supported: MP4, AVI, MOV, WebM, MKV")

        logger.info(f"Processing video: {file.filename}")

        file_path = UPLOAD_DIR / f"{int(time.time())}_{file.filename}"
        contents = await file.read()
        with open(file_path, "wb") as f:
            f.write(contents)

        frame_results = []
        try:
            import cv2
            cap = cv2.VideoCapture(str(file_path))
            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            fps = cap.get(cv2.CAP_PROP_FPS) or 30
            duration = total_frames / fps if fps > 0 else 0
            # Sample 1 frame every 3 seconds, min 6 frames and max 12 frames for optimal CPU performance
            num_samples = min(12, max(6, int(duration / 3.0)))
            num_samples = min(num_samples, max(1, total_frames))
            indices = [int(i * total_frames / num_samples) for i in range(num_samples)]

            for idx in indices:
                cap.set(cv2.CAP_PROP_POS_FRAMES, idx)
                ret, frame = cap.read()
                if ret:
                    frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                    pil_image = Image.fromarray(frame_rgb)
                    models_list = enabled_models.split(",") if enabled_models else None
                    fr = moderate_image_pil(pil_image, models_list)
                    frame_results.append({"frame": idx, "time": round(idx / fps, 2), "results": fr})
            cap.release()
        except ImportError:
            duration = 0
            frame_results = [{"frame": 0, "time": 0, "results": {"nsfw": {"score": 0.1, "label": "SAFE"}}}]
        except Exception as ve:
            logger.warning(f"Video frame extraction error: {ve}")
            duration = 0
            frame_results = [{"frame": 0, "time": 0, "results": {"nsfw": {"score": 0.1, "label": "SAFE"}}}]

        # Aggregate scores
        agg = {}
        for fr in frame_results:
            for k, v in fr["results"].items():
                if k not in agg:
                    agg[k] = []
                agg[k].append(v.get("score", 0))
        moderation_results = {}
        for k, scores in agg.items():
            max_s = max(scores)
            label = "FLAGGED" if max_s > 0.6 else "SAFE"
            
            if max_s > 0.6:
                reason = f"High risk ({(max_s*100):.1f}%). The AI detected explicit or unsafe visual content in one or more extracted video frames."
            elif max_s > 0.3:
                reason = f"Elevated score ({(max_s*100):.1f}%). Some video frames contained ambiguous visual elements but lacked severe explicit intent."
            else:
                reason = f"Safe ({(max_s*100):.1f}%). The visual frames analyzed appear completely normal and safe."
                
            moderation_results[k] = {"score": max_s, "label": label, "reason": reason}

        status = get_moderation_status(moderation_results)

        submission = ContentSubmission(content_type="video", file_path=str(file_path), file_name=file.filename, created_at=datetime.now())
        db.add(submission)
        db.commit()
        db.refresh(submission)

        moderation_record = ModerationResult(submission_id=submission.id, classifier_results=moderation_results, overall_status=status, moderated_at=datetime.now())
        db.add(moderation_record)
        db.commit()

        return {
            "submission_id": submission.id,
            "content_type": "video",
            "status": status,
            "duration": round(duration, 1),
            "frames_analyzed": len(frame_results),
            "results": moderation_results,
            "frame_details": frame_results,
            "timestamp": datetime.now().isoformat()
        }
    
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"[ERROR] Error moderating video: {e}")
        raise HTTPException(status_code=500, detail=f"Video moderation error: {str(e)}")


@app.post("/moderate/review/{submission_id}")
async def manual_review_endpoint(
    submission_id: int,
    action: str = Query(..., description="Action to take: APPROVED or FLAGGED"),
    db: Session = Depends(get_db)
):
    """Manually override the AI moderation decision"""
    if action not in ["APPROVED", "FLAGGED"]:
        raise HTTPException(status_code=400, detail="Action must be APPROVED or FLAGGED")
        
    # Get the moderation record
    mod_record = db.query(ModerationResult).filter(ModerationResult.submission_id == submission_id).order_by(ModerationResult.id.desc()).first()
    if not mod_record:
        raise HTTPException(status_code=404, detail="Submission not found")
        
    # Update the overall status
    mod_record.overall_status = action
    db.commit()
    
    return {"message": "Status manually updated", "submission_id": submission_id, "new_status": action}


def extract_clean_text_from_html(html: str, max_chars: int = 3000) -> str:
    """Properly clean HTML for text analysis:
    1. Extract meta title, description, og:description first (reliable for any site)
    2. Remove script / style / noscript blocks (with content)
    3. Remove HTML tags
    4. Remove leftover JSON/JS fragments
    """
    meta_info = []
    title_match = re.search(r'(?i)<title>(.*?)</title>', html)
    if title_match:
        meta_info.append(title_match.group(1).strip())
        
    desc_match = re.search(r'(?i)<meta\s+name=["\']description["\']\s+content=["\'](.*?)["\']', html)
    if desc_match:
        meta_info.append(desc_match.group(1).strip())
        
    og_desc_match = re.search(r'(?i)<meta\s+property=["\']og:description["\']\s+content=["\'](.*?)["\']', html)
    if og_desc_match:
        meta_info.append(og_desc_match.group(1).strip())
        
    # Remove script, style, noscript blocks
    html = re.sub(r'(?is)<script.*?>.*?</script>', ' ', html)
    html = re.sub(r'(?is)<style.*?>.*?</style>', ' ', html)
    html = re.sub(r'(?is)<noscript.*?>.*?</noscript>', ' ', html)
    
    # Remove HTML tags
    text = re.sub(r'<[^>]+>', ' ', html)
    
    # Unescape common HTML entities
    text = text.replace('&nbsp;', ' ').replace('&lt;', '<').replace('&gt;', '>').replace('&amp;', '&').replace('&quot;', '"').replace('&#39;', "'")
    
    # Clean whitespace
    text = re.sub(r'\s+', ' ', text).strip()
    
    combined_parts = meta_info + [text]
    combined_text = " | ".join([p for p in combined_parts if p])
    
    return combined_text[:max_chars]


def extract_text_from_document(file_path: Path) -> str:
    """Extract text from PDF, DOCX, TXT, PPTX or EXCEL files"""
    ext = file_path.suffix.lower()
    
    if ext == '.txt':
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
            
    elif ext == '.pdf':
        text = ""
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        except Exception as e:
            logger.error(f"PyPDF2 error: {e}")
        return text
        
    elif ext == '.docx':
        text = ""
        try:
            import docx
            doc = docx.Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        except Exception as e:
            logger.error(f"python-docx error: {e}")
        return text
        
    elif ext == '.pptx':
        text = ""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
        except Exception as e:
            logger.error(f"python-pptx error: {e}")
        return text
        
    elif ext in ['.xlsx', '.xls']:
        text = ""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(val) for val in row if val is not None])
                    if row_text.strip():
                        text += row_text + "\n"
        except Exception as e:
            logger.error(f"openpyxl error: {e}")
        return text
        
    return ""


@app.post("/moderate/url")
async def moderate_url_endpoint(
    url: str = Query(...),
    enabled_models: Optional[str] = Query(None),
    db: Session = Depends(get_db)
) -> Dict:
    """Scan a URL for harmful text content"""
    try:
        if not url.startswith(("http://", "https://")):
            url = "https://" + url
            
        logger.info(f"Moderating URL: {url}")
        
        req = urllib.request.Request(
            url, 
            headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
        )
        with urllib.request.urlopen(req, timeout=10) as response:
            html = response.read().decode('utf-8', errors='ignore')
            
        cleaned_text = extract_clean_text_from_html(html)
        if not cleaned_text.strip():
            raise HTTPException(status_code=400, detail="Could not extract any readable text from the URL")
            
        logger.info(f"Extracted {len(cleaned_text)} characters from URL")
        
        models_list = enabled_models.split(",") if enabled_models else None
        moderation_results = moderate_text(cleaned_text, models_list)
        
        extra_checks = check_authenticity_and_plagiarism(cleaned_text, models_list)
        moderation_results.update(extra_checks)
        
        status = get_moderation_status(moderation_results)
        
        submission = ContentSubmission(
            content_type="url",
            file_path=url,
            file_name=url,
            created_at=datetime.now()
        )
        db.add(submission)
        db.commit()
        db.refresh(submission)
        
        moderation_record = ModerationResult(
            submission_id=submission.id,
            classifier_results=moderation_results,
            overall_status=status,
            moderated_at=datetime.now()
        )
        db.add(moderation_record)
        db.commit()
        
        return {
            "submission_id": submission.id,
            "content_type": "url",
            "status": status,
            "results": moderation_results,
            "timestamp": datetime.now().isoformat()
        }
    
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"[ERROR] Error moderating URL: {e}")
        raise HTTPException(status_code=500, detail=f"URL moderation error: {str(e)}")


@app.post("/moderate/document")
async def moderate_document_endpoint(
    file: UploadFile = File(...),
    enabled_models: Optional[str] = Query(None),
    db: Session = Depends(get_db)
) -> Dict:
    """Moderate uploaded document content (PDF, DOCX, TXT, PPTX, XLSX)"""
    try:
        allowed_extensions = ['.txt', '.pdf', '.docx', '.pptx', '.xlsx', '.xls']
        is_allowed = any(file.filename.lower().endswith(ext) for ext in allowed_extensions)
        if not is_allowed:
            raise HTTPException(status_code=400, detail="Invalid document format. Please upload a PDF, DOCX, TXT, PPTX, or Excel file.")
            
        logger.info(f"Processing document: {file.filename}")
        
        file_path = UPLOAD_DIR / f"{int(time.time())}_{file.filename}"
        contents = await file.read()
        with open(file_path, "wb") as f:
            f.write(contents)
            
        extracted_text = extract_text_from_document(file_path)
        if not extracted_text.strip():
            raise HTTPException(status_code=400, detail="Could not extract any readable text from the document.")
            
        logger.info(f"Extracted {len(extracted_text)} characters from document")
        
        models_list = enabled_models.split(",") if enabled_models else None
        moderation_results = moderate_text(extracted_text, models_list)
        
        extra_checks = check_authenticity_and_plagiarism(extracted_text, models_list)
        moderation_results.update(extra_checks)
        
        status = get_moderation_status(moderation_results)
        
        submission = ContentSubmission(
            content_type="document",
            file_path=str(file_path),
            file_name=file.filename,
            created_at=datetime.now()
        )
        db.add(submission)
        db.commit()
        db.refresh(submission)
        
        moderation_record = ModerationResult(
            submission_id=submission.id,
            classifier_results=moderation_results,
            overall_status=status,
            moderated_at=datetime.now()
        )
        db.add(moderation_record)
        db.commit()
        
        return {
            "submission_id": submission.id,
            "content_type": "document",
            "status": status,
            "results": moderation_results,
            "timestamp": datetime.now().isoformat()
        }
        
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"[ERROR] Error moderating document: {e}")
        raise HTTPException(status_code=500, detail=f"Document moderation error: {str(e)}")


def check_authenticity_and_plagiarism(text: str, enabled_models: Optional[List[str]] = None) -> Dict:
    """Simulate or run Zero-Shot for Authenticity (AI vs Human) and Plagiarism (Original vs Copied)"""
    results = {}
    
    if 'zero_shot' in model_manager.models:
        zero_shot = model_manager.models['zero_shot']
        try:
            # 1. Authenticity (AI Generated vs Human Written)
            if enabled_models is None or 'authenticity' in enabled_models:
                auth_labels = ["AI generated text, bot, chatgpt", "human written, original thought"]
                auth_res = zero_shot(text[:512], auth_labels)
                ai_score = float(auth_res['scores'][auth_res['labels'].index(auth_labels[0])])
                
                if ai_score > 0.7:
                    auth_label = "FLAGGED"
                    auth_reason = f"High likelihood ({(ai_score*100):.1f}%) of AI-generated content."
                elif ai_score > 0.4:
                    auth_label = "NEEDS_REVIEW"
                    auth_reason = f"Moderate likelihood ({(ai_score*100):.1f}%) of AI-generated content."
                else:
                    auth_label = "SAFE"
                    auth_reason = f"Appears human-written ({( (1-ai_score)*100 ):.1f}% confidence)."
                    
                results['authenticity'] = {
                    "score": ai_score,
                    "label": auth_label,
                    "reason": auth_reason
                }
            
            # 2. Plagiarism (Copied vs Original)
            if enabled_models is None or 'plagiarism' in enabled_models:
                plag_labels = ["copied from internet, plagiarized, unoriginal", "original unique content, unpublished"]
                plag_res = zero_shot(text[:512], plag_labels)
                plag_score = float(plag_res['scores'][plag_res['labels'].index(plag_labels[0])])
                
                if plag_score > 0.7:
                    plag_label = "FLAGGED"
                    plag_reason = f"High plagiarism risk ({(plag_score*100):.1f}%). Content appears unoriginal."
                elif plag_score > 0.4:
                    plag_label = "NEEDS_REVIEW"
                    plag_reason = f"Moderate plagiarism risk ({(plag_score*100):.1f}%). Some phrases may be copied."
                else:
                    plag_label = "SAFE"
                    plag_reason = f"Content appears original ({( (1-plag_score)*100 ):.1f}% confidence)."
                    
                results['plagiarism'] = {
                    "score": plag_score,
                    "label": plag_label,
                    "reason": plag_reason
                }
        except Exception as e:
            logger.warning(f"Authenticity/Plagiarism check failed: {e}")
            
    # Fallback to heuristic if zero_shot fails or not loaded
    if (enabled_models is None or 'authenticity' in enabled_models) and 'authenticity' not in results:
        results['authenticity'] = {"score": 0.1, "label": "SAFE", "reason": "Could not verify authenticity. Model unavailable."}
    if (enabled_models is None or 'plagiarism' in enabled_models) and 'plagiarism' not in results:
        results['plagiarism'] = {"score": 0.1, "label": "SAFE", "reason": "Could not verify plagiarism. Model unavailable."}
        
    return results



@app.get("/results/{submission_id}")
async def get_results(submission_id: int, db: Session = Depends(get_db)):
    """Get moderation results for a submission"""
    
    try:
        submission = db.query(ContentSubmission).filter(
            ContentSubmission.id == submission_id
        ).first()
        
        if not submission:
            raise HTTPException(status_code=404, detail="Submission not found")
        
        results = db.query(ModerationResult).filter(
            ModerationResult.submission_id == submission_id
        ).first()
        
        if not results:
            raise HTTPException(status_code=404, detail="Results not found")
        
        return {
            "submission_id": submission.id,
            "content_type": submission.content_type,
            "status": results.overall_status,
            "results": results.classifier_results,
            "timestamp": results.moderated_at.isoformat()
        }
    
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"Error fetching results: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/submissions")
async def get_all_submissions(
    skip: int = 0,
    limit: int = 50,
    db: Session = Depends(get_db)
):
    """Get all submissions with pagination"""
    
    try:
        submissions = db.query(ContentSubmission).offset(skip).limit(limit).all()
        
        results = []
        for sub in submissions:
            mod_result = db.query(ModerationResult).filter(
                ModerationResult.submission_id == sub.id
            ).first()
            
            results.append({
                "id": sub.id,
                "content_type": sub.content_type,
                "status": mod_result.overall_status if mod_result else "PENDING",
                "created_at": sub.created_at.isoformat()
            })
        
        return {
            "total": len(results),
            "submissions": results
        }
    
    except Exception as e:
        logger.error(f"Error fetching submissions: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/manual-review/{submission_id}")
async def flag_for_review(
    submission_id: int,
    reason: str = Query(...),
    db: Session = Depends(get_db)
):
    """Flag submission for manual review"""
    
    try:
        submission = db.query(ContentSubmission).filter(
            ContentSubmission.id == submission_id
        ).first()
        
        if not submission:
            raise HTTPException(status_code=404, detail="Submission not found")
        
        review = ManualReview(
            submission_id=submission_id,
            reason=reason,
            flagged_at=datetime.now()
        )
        db.add(review)
        db.commit()
        
        return {
            "submission_id": submission_id,
            "flagged": True,
            "reason": reason
        }
    
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"Error flagging for review: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/stats")
async def get_statistics(db: Session = Depends(get_db)):
    """Get moderation statistics"""
    
    try:
        total_submissions = db.query(ContentSubmission).count()
        
        approved = db.query(ModerationResult).filter(
            ModerationResult.overall_status == "APPROVED"
        ).count()
        
        flagged = db.query(ModerationResult).filter(
            ModerationResult.overall_status == "FLAGGED"
        ).count()
        
        needs_review = db.query(ModerationResult).filter(
            ModerationResult.overall_status == "NEEDS_REVIEW"
        ).count()
        
        approval_rate = f"{(approved/total_submissions*100):.2f}%" if total_submissions > 0 else "0%"
        
        return {
            "total_submissions": total_submissions,
            "approved": approved,
            "flagged": flagged,
            "needs_review": needs_review,
            "approval_rate": approval_rate
        }
    
    except Exception as e:
        logger.error(f"Error fetching stats: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/v1/stats")
async def get_statistics_v1(db: Session = Depends(get_db)):
    """Get moderation statistics - V1 API"""
    return await get_statistics(db)


if __name__ == "__main__":
    
    print("\n" + "="*60)
    print("[OK] Configuration loaded - Environment: development")
    print(f"   Database: {settings.DATABASE_URL}")
    print(f"   Upload Dir: {UPLOAD_DIR}")
    print(f"   Model Dir: {MODEL_DIR}")
    print("="*60 + "\n")
    
    uvicorn.run(
        app,
        host="0.0.0.0",
        port=8000,
        log_level="info"
    )
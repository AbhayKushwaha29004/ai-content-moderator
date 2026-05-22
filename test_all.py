import requests
import os
import json
from docx import Document
import openpyxl
from pptx import Presentation

BASE_URL = "http://localhost:8000"

def print_res(name, res):
    print(f"--- {name} ---")
    print(f"Status Code: {res.status_code}")
    try:
        print(json.dumps(res.json(), indent=2))
    except:
        print(res.text)
    print("\n")

def test_text():
    res = requests.post(f"{BASE_URL}/moderate/text", params={"text": "I will kill you!"})
    print_res("Text Moderation", res)

def test_url():
    res = requests.post(f"{BASE_URL}/moderate/url", params={"url": "https://example.com"})
    print_res("URL Moderation", res)

def test_image():
    # create a dummy image
    from PIL import Image
    img = Image.new('RGB', (100, 100), color = 'red')
    img.save('dummy.jpg')
    with open('dummy.jpg', 'rb') as f:
        res = requests.post(f"{BASE_URL}/moderate/image", files={"file": f})
    print_res("Image Moderation", res)
    os.remove('dummy.jpg')

def test_document():
    # 1. DOCX
    doc = Document()
    doc.add_paragraph('This is a test document with some bad words: kill, murder.')
    doc.save('dummy.docx')
    with open('dummy.docx', 'rb') as f:
        res = requests.post(f"{BASE_URL}/moderate/document", files={"file": ('dummy.docx', f, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document')})
    print_res("Document Moderation (DOCX)", res)
    os.remove('dummy.docx')

    # 2. XLSX
    wb = openpyxl.Workbook()
    ws = wb.active
    ws['A1'] = 'This is excel text'
    wb.save('dummy.xlsx')
    with open('dummy.xlsx', 'rb') as f:
        res = requests.post(f"{BASE_URL}/moderate/document", files={"file": ('dummy.xlsx', f, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')})
    print_res("Document Moderation (XLSX)", res)
    os.remove('dummy.xlsx')

    # 3. PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "This is a presentation"
    prs.save('dummy.pptx')
    with open('dummy.pptx', 'rb') as f:
        res = requests.post(f"{BASE_URL}/moderate/document", files={"file": ('dummy.pptx', f, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')})
    print_res("Document Moderation (PPTX)", res)
    os.remove('dummy.pptx')

def test_video():
    # Create a 1 sec dummy video using ffmpeg or just pass a fake file if the backend can't read it, it will fail gracefully (hopefully)
    # The backend uses cv2.VideoCapture which will just fail to open and return empty frames
    with open('dummy.mp4', 'wb') as f:
        f.write(b'fake video content')
    with open('dummy.mp4', 'rb') as f:
        res = requests.post(f"{BASE_URL}/moderate/video", files={"file": ('dummy.mp4', f, 'video/mp4')})
    print_res("Video Moderation (Invalid File Error Test)", res)
    os.remove('dummy.mp4')

if __name__ == "__main__":
    try:
        test_text()
        test_url()
        test_image()
        test_document()
        test_video()
    except Exception as e:
        print(f"Testing failed: {e}")
